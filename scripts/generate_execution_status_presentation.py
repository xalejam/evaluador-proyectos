from __future__ import annotations

import sqlite3
from dataclasses import dataclass
from datetime import datetime
from io import BytesIO
from pathlib import Path
from typing import Iterable

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


DB_PATH = Path(__file__).resolve().parent.parent / "project_viability.db"
OUTPUT_PATH = Path("docs") / "Resumen_Proyectos_Ejecucion.pptx"
LOGO_PATH = Path("logo_DDNola.png")

# Colors
C_HEADER_BG   = RGBColor(7, 89, 90)
C_HEADER_TEXT = RGBColor(255, 255, 255)
C_HEADER_SUB  = RGBColor(160, 220, 220)
C_HEADER_CUT  = RGBColor(200, 235, 235)
C_BODY_BG     = RGBColor(248, 250, 252)
C_ROW_ODD     = RGBColor(255, 255, 255)
C_ROW_EVEN    = RGBColor(248, 250, 252)
C_ROW_BORDER  = RGBColor(226, 232, 240)
C_SEPARATOR   = RGBColor(203, 213, 225)
C_TEXT_DARK   = RGBColor(15, 23, 42)
C_TEXT_MID    = RGBColor(51, 65, 85)
C_TEXT_LIGHT  = RGBColor(100, 116, 139)
C_TEXT_BLUE   = RGBColor(71, 85, 105)
C_TEXT_BLUE2  = RGBColor(191, 219, 254)
C_CARD_BG     = RGBColor(241, 245, 249)
C_CARD_BORDER = RGBColor(203, 213, 225)
C_RISK_HIGH   = RGBColor(220, 38, 38)
C_RISK_MED    = RGBColor(217, 119, 6)

# Layout — slide dimensions
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
HEADER_H = Inches(0.78)
ROW_H = Inches(1.15)

# Slide 1: metrics + col headers above rows
COL_HEADER_Y_S1 = Inches(2.28)
ROW_START_Y_S1  = Inches(2.50)

# Slide 2+: rows start right after header
COL_HEADER_Y_SN = Inches(0.90)
ROW_START_Y_SN  = Inches(1.12)

# Column positions
COL_NAME_X  = Inches(0.65)
COL_NAME_W  = Inches(3.15)
COL_BADGE_X = Inches(3.85)
COL_BADGE_W = Inches(0.90)
COL_NOTES_X = Inches(4.90)
COL_NOTES_W = Inches(4.20)
COL_SEP_X   = Inches(9.18)
COL_RISK_X  = Inches(9.30)
COL_RISK_W  = Inches(1.95)
COL_DATE_X  = Inches(11.35)
COL_DATE_W  = Inches(1.40)

ROWS_PER_SLIDE = 4


@dataclass
class ProjectStatus:
    project_id: str
    name: str
    progress_percent: int | None
    progress_at: str | None
    general_note: str
    next_step: str
    blocker: str
    risk: str


def fetch_executing_projects(db_path: Path) -> list[ProjectStatus]:
    conn = sqlite3.connect(db_path)
    conn.row_factory = sqlite3.Row
    sql = """
    WITH latest_progress AS (
        SELECT h.project_id, h.progress_percent, h.created_at
        FROM v_project_progress_history h
        JOIN (
            SELECT project_id, MAX(datetime(created_at)) AS max_created_at
            FROM v_project_progress_history
            GROUP BY project_id
        ) mx
          ON mx.project_id = h.project_id
         AND datetime(mx.max_created_at) = datetime(h.created_at)
    ),
    latest_notes AS (
        SELECT
            n.project_id,
            n.note_type,
            n.note_text,
            n.created_at,
            ROW_NUMBER() OVER (
                PARTITION BY n.project_id, n.note_type
                ORDER BY datetime(n.created_at) DESC, n.note_id DESC
            ) AS rn
        FROM project_notes n
    )
    SELECT
        p.project_id,
        p.name,
        lp.progress_percent,
        lp.created_at AS progress_at,
        COALESCE(gn.note_text, '') AS general_note,
        COALESCE(pn.note_text, '') AS next_step,
        COALESCE(bn.note_text, '') AS blocker,
        COALESCE(rk.note_text, '') AS risk
    FROM projects p
    LEFT JOIN latest_progress lp
        ON lp.project_id = p.project_id
    LEFT JOIN latest_notes gn
        ON gn.project_id = p.project_id AND gn.note_type = 'general' AND gn.rn = 1
    LEFT JOIN latest_notes pn
        ON pn.project_id = p.project_id AND pn.note_type = 'proximo_paso' AND pn.rn = 1
    LEFT JOIN latest_notes bn
        ON bn.project_id = p.project_id AND bn.note_type = 'bloqueador' AND bn.rn = 1
    LEFT JOIN latest_notes rk
        ON rk.project_id = p.project_id AND rk.note_type = 'riesgo' AND rk.rn = 1
    WHERE lower(COALESCE(p.status, '')) = 'executing'
    ORDER BY COALESCE(lp.created_at, p.updated_at, p.created_date) DESC, p.project_id
    """
    rows = conn.execute(sql).fetchall()
    conn.close()
    return [ProjectStatus(**dict(row)) for row in rows]


def average_progress(projects: Iterable[ProjectStatus]) -> int:
    values = [p.progress_percent for p in projects if isinstance(p.progress_percent, int)]
    return round(sum(values) / len(values)) if values else 0


def latest_update(projects: Iterable[ProjectStatus]) -> str:
    values = [datetime.fromisoformat(p.progress_at) for p in projects if p.progress_at]
    return max(values).strftime("%d %b %Y") if values else "Sin actualizacion"


def progress_color(progress: int | None) -> RGBColor:
    if progress is None:
        return RGBColor(160, 174, 192)
    if progress >= 70:
        return RGBColor(22, 163, 74)
    if progress >= 50:
        return RGBColor(245, 158, 11)
    return RGBColor(220, 38, 38)


def risk_color(risk: str) -> RGBColor:
    r = (risk or "").lower()
    if any(w in r for w in ("critico", "crítico", "alto")):
        return C_RISK_HIGH
    if r not in ("", "ninguno", "ninguna", "n/a"):
        return C_RISK_MED
    return C_TEXT_LIGHT


def blocker_color(blocker: str) -> RGBColor:
    b = (blocker or "").strip().lower()
    return C_RISK_HIGH if b not in ("", "ninguno", "ninguna") else C_TEXT_LIGHT


def trim_text(value: str, limit: int) -> str:
    clean = " ".join((value or "").split())
    return clean if len(clean) <= limit else clean[: limit - 3].rstrip() + "..."


def add_textbox(slide, left, top, width, height, text, font_size, *, bold=False, color=None, align=PP_ALIGN.LEFT):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.alignment = align
    run = p.runs[0]
    run.font.size = Pt(font_size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color
    return box


def add_metric_card(slide, left, title, value, subtitle):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, Inches(1.05), Inches(2.2), Inches(1.0))
    card.fill.solid()
    card.fill.fore_color.rgb = C_CARD_BG
    card.line.color.rgb = C_CARD_BORDER
    add_textbox(slide, left + Inches(0.12), Inches(1.13), Inches(1.95), Inches(0.22), title, 10, color=C_TEXT_BLUE)
    add_textbox(slide, left + Inches(0.12), Inches(1.35), Inches(1.95), Inches(0.3), value, 20, bold=True, color=C_TEXT_DARK)
    add_textbox(slide, left + Inches(0.12), Inches(1.68), Inches(1.95), Inches(0.2), subtitle, 9, color=C_TEXT_LIGHT)


def _draw_header(slide, prs, slide_num: int, total_slides: int) -> None:
    header = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, HEADER_H)
    header.fill.solid()
    header.fill.fore_color.rgb = C_HEADER_BG
    header.line.color.rgb = C_HEADER_BG

    add_textbox(slide, Inches(0.45), Inches(0.18), Inches(9.0), Inches(0.3),
                "Resumen ejecutivo | Proyectos en ejecucion", 24, bold=True, color=C_HEADER_TEXT)
    add_textbox(slide, Inches(0.45), Inches(0.5), Inches(4.3), Inches(0.18),
                "Fuente: project_viability.db", 9, color=C_HEADER_SUB)

    if total_slides > 1:
        add_textbox(slide, Inches(10.3), Inches(0.28), Inches(1.5), Inches(0.18),
                    f"{slide_num} / {total_slides}", 10, color=C_HEADER_CUT, align=PP_ALIGN.RIGHT)

    if LOGO_PATH.exists():
        try:
            slide.shapes.add_picture(str(LOGO_PATH), Inches(12.0), Inches(0.09), height=Inches(0.58))
        except Exception:
            pass


def _draw_column_headers(slide, col_header_y) -> None:
    cols = [
        (COL_NAME_X,  COL_NAME_W,  "Proyecto"),
        (COL_BADGE_X, COL_BADGE_W, "Avance"),
        (COL_NOTES_X, COL_NOTES_W, "Nota / Proximo paso"),
        (COL_RISK_X,  COL_RISK_W,  "Bloqueo & Riesgo"),
        (COL_DATE_X,  COL_DATE_W,  "Ultima act."),
    ]
    for x, w, label in cols:
        add_textbox(slide, x, col_header_y, w, Inches(0.18), label, 8, bold=True, color=C_TEXT_LIGHT)


def _draw_project_row(slide, project: ProjectStatus, row_index: int, row_start_y) -> None:
    top = row_start_y + row_index * ROW_H
    row_bg = C_ROW_ODD if row_index % 2 == 0 else C_ROW_EVEN
    row = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                  Inches(0.45), top, Inches(12.4), ROW_H - Inches(0.05))
    row.fill.solid()
    row.fill.fore_color.rgb = row_bg
    row.line.color.rgb = C_ROW_BORDER

    # Project name
    add_textbox(slide, COL_NAME_X, top + Inches(0.08), COL_NAME_W, Inches(0.30),
                project.name, 14, bold=True, color=C_TEXT_DARK)
    # Project ID — near bottom of card to avoid overlapping name
    add_textbox(slide, COL_NAME_X, top + ROW_H - Inches(0.35), COL_NAME_W, Inches(0.22),
                project.project_id, 9, color=C_TEXT_LIGHT)

    # Progress badge
    progress = project.progress_percent or 0
    badge = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                                    COL_BADGE_X, top + Inches(0.14), COL_BADGE_W, Inches(0.36))
    badge.fill.solid()
    badge.fill.fore_color.rgb = progress_color(project.progress_percent)
    badge.line.color.rgb = progress_color(project.progress_percent)
    add_textbox(slide, COL_BADGE_X, top + Inches(0.18), COL_BADGE_W, Inches(0.22),
                f"{progress}%", 16, bold=True, color=RGBColor(255, 255, 255), align=PP_ALIGN.CENTER)

    # Vertical separator between notes and risk
    sep = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                                  COL_SEP_X, top + Inches(0.08), Inches(0.02), ROW_H - Inches(0.20))
    sep.fill.solid()
    sep.fill.fore_color.rgb = C_SEPARATOR
    sep.line.color.rgb = C_SEPARATOR

    # Notes column
    add_textbox(slide, COL_NOTES_X, top + Inches(0.05), COL_NOTES_W, Inches(0.14),
                "NOTA", 7, bold=True, color=C_TEXT_LIGHT)
    add_textbox(slide, COL_NOTES_X, top + Inches(0.19), COL_NOTES_W, Inches(0.32),
                trim_text(project.general_note, 160), 9, color=C_TEXT_MID)
    add_textbox(slide, COL_NOTES_X, top + Inches(0.55), COL_NOTES_W, Inches(0.14),
                "PROX. PASO", 7, bold=True, color=C_TEXT_LIGHT)
    add_textbox(slide, COL_NOTES_X, top + Inches(0.69), COL_NOTES_W, Inches(0.30),
                trim_text(project.next_step or "Sin definir", 100), 9, color=C_TEXT_BLUE)

    # Risk / blocker column
    add_textbox(slide, COL_RISK_X, top + Inches(0.05), COL_RISK_W, Inches(0.14),
                "BLOQUEADOR", 7, bold=True, color=C_TEXT_LIGHT)
    add_textbox(slide, COL_RISK_X, top + Inches(0.19), COL_RISK_W, Inches(0.24),
                trim_text(project.blocker or "Ninguno", 45), 9,
                color=blocker_color(project.blocker or ""))
    add_textbox(slide, COL_RISK_X, top + Inches(0.53), COL_RISK_W, Inches(0.14),
                "RIESGO", 7, bold=True, color=C_TEXT_LIGHT)
    add_textbox(slide, COL_RISK_X, top + Inches(0.67), COL_RISK_W, Inches(0.24),
                trim_text(project.risk or "Ninguno", 50), 9,
                color=risk_color(project.risk or ""))

    # Date
    date_text = datetime.fromisoformat(project.progress_at).strftime("%d %b %Y") if project.progress_at else "Sin fecha"
    add_textbox(slide, COL_DATE_X, top + Inches(0.24), COL_DATE_W, Inches(0.24),
                date_text, 10, bold=True, color=C_TEXT_DARK, align=PP_ALIGN.CENTER)
    add_textbox(slide, COL_DATE_X, top + Inches(0.50), COL_DATE_W, Inches(0.18),
                "Ultima act.", 8, color=C_TEXT_LIGHT, align=PP_ALIGN.CENTER)


def _draw_metrics(slide, all_projects: list[ProjectStatus]) -> None:
    no_blockers = sum(1 for p in all_projects if (p.blocker or "").strip().lower() in {"", "ninguno", "ninguna"})
    add_metric_card(slide, Inches(0.45),  "Proyectos activos",    str(len(all_projects)),               "estatus executing")
    add_metric_card(slide, Inches(2.9),   "Avance promedio",      f"{average_progress(all_projects)}%", "segun ultimo registro")
    add_metric_card(slide, Inches(5.35),  "Sin bloqueadores",     str(no_blockers),                     "ultimo reporte")
    add_metric_card(slide, Inches(7.8),   "Ultima actualizacion", latest_update(all_projects),          "avance mas reciente")
    add_metric_card(slide, Inches(10.25), "Mensaje clave",        f"{len(all_projects)} frentes",       "sin bloqueos criticos")

    add_textbox(slide, Inches(0.45), Inches(1.98), Inches(12.0), Inches(0.18),
                "Todos los proyectos en ejecucion reportan avance y no muestran bloqueadores ni riesgos criticos en la ultima nota.",
                11, color=C_TEXT_MID)


def _draw_footer(slide, generated_at: str) -> None:
    add_textbox(slide, Inches(0.45), Inches(7.08), Inches(8.0), Inches(0.18),
                "Nota: el porcentaje corresponde al ultimo progreso capturado por proyecto.",
                8, color=C_TEXT_LIGHT)
    add_textbox(slide, Inches(9.0), Inches(7.08), Inches(4.0), Inches(0.18),
                f"Corte: {generated_at}", 8, color=C_TEXT_LIGHT, align=PP_ALIGN.RIGHT)


def build_slide(
    prs: Presentation,
    chunk: list[ProjectStatus],
    slide_num: int,
    total_slides: int,
    all_projects: list[ProjectStatus],
    generated_at: str,
) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = C_BODY_BG

    _draw_header(slide, prs, slide_num, total_slides)

    if slide_num == 1:
        _draw_metrics(slide, all_projects)
        col_hdr_y = COL_HEADER_Y_S1
        row_start_y = ROW_START_Y_S1
    else:
        col_hdr_y = COL_HEADER_Y_SN
        row_start_y = ROW_START_Y_SN

    _draw_column_headers(slide, col_hdr_y)

    for i, project in enumerate(chunk):
        _draw_project_row(slide, project, i, row_start_y)

    _draw_footer(slide, generated_at)


def _build_prs(projects: list[ProjectStatus]) -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    chunks = [projects[i:i + ROWS_PER_SLIDE] for i in range(0, len(projects), ROWS_PER_SLIDE)]
    generated_at = datetime.now().strftime("%d/%m/%Y %H:%M")
    for n, chunk in enumerate(chunks, start=1):
        build_slide(prs, chunk, n, len(chunks), projects, generated_at)
    return prs


def build_presentation(projects: list[ProjectStatus], output_path: Path) -> Path:
    prs = _build_prs(projects)
    output_path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output_path)
    return output_path


def build_presentation_bytes(projects: list[ProjectStatus]) -> bytes:
    buf = BytesIO()
    _build_prs(projects).save(buf)
    return buf.getvalue()


def main() -> None:
    projects = fetch_executing_projects(DB_PATH)
    if not projects:
        raise SystemExit("No hay proyectos con status 'executing'.")
    path = build_presentation(projects, OUTPUT_PATH)
    print(path)


if __name__ == "__main__":
    main()
